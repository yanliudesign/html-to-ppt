#!/usr/bin/env python3
"""
HTML slide deck  →  editable PowerPoint (.pptx).

Every text block, image, colored box and border becomes a native PowerPoint
object you can select, edit, restyle, move, or delete individually.

Only visuals that CSS can express but PowerPoint cannot (gradients,
gradient-clipped text, inline SVG radar charts, radial art blobs) are
captured as a discrete picture object so they stay editable-as-picture
(replaceable, movable, resizable) instead of being flattened into the slide.

Usage:
    python3 tools/html_to_pptx_editable.py <input.html> <output.pptx>
"""

from __future__ import annotations

import io
import os
import re
import sys
import json
import tempfile
from pathlib import Path
from typing import Any

from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Emu, Pt, Length
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------
SLIDE_W_PX = 1920
SLIDE_H_PX = 1080
EMU_PER_PX = 9525               # 96 dpi
DPR = 1                          # screenshot device pixel ratio (1 = fast; bump to 2 for crisper bitmaps)


def px_emu(v: float) -> int:
    return int(round(v * EMU_PER_PX))


# ---------------------------------------------------------------------------
# Color parsing
# ---------------------------------------------------------------------------
_RGB_RE = re.compile(
    r"rgba?\(\s*([\d.]+)\s*,\s*([\d.]+)\s*,\s*([\d.]+)\s*(?:,\s*([\d.]+)\s*)?\)"
)
_HEX_RE = re.compile(r"^#([0-9a-fA-F]{3,8})$")


def parse_color(css: str | None):
    """Return (r,g,b,a) with a in [0,1] or None if fully transparent/unknown."""
    if not css:
        return None
    css = css.strip()
    if css in ("transparent", "none", "currentcolor"):
        return None
    m = _RGB_RE.match(css)
    if m:
        r, g, b = int(float(m.group(1))), int(float(m.group(2))), int(float(m.group(3)))
        a = float(m.group(4)) if m.group(4) is not None else 1.0
        if a <= 0:
            return None
        return (r, g, b, a)
    m = _HEX_RE.match(css)
    if m:
        h = m.group(1)
        if len(h) == 3:
            r, g, b = (int(h[i] * 2, 16) for i in range(3))
            return (r, g, b, 1.0)
        if len(h) == 6:
            return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16), 1.0)
        if len(h) == 8:
            return (
                int(h[0:2], 16),
                int(h[2:4], 16),
                int(h[4:6], 16),
                int(h[6:8], 16) / 255,
            )
    return None


def to_rgb(rgba):
    if not rgba:
        return None
    return RGBColor(rgba[0], rgba[1], rgba[2])


# ---------------------------------------------------------------------------
# Extractor JS — runs inside the page, returns an array of visual items per slide.
# ---------------------------------------------------------------------------
EXTRACT_JS = r"""
(slideIdx) => {
  const SKIP_TAGS = new Set(['SCRIPT','STYLE','META','LINK','HEAD','NOSCRIPT','TEMPLATE','BR']);
  const SKIP_CLASSES = ['nav','edit-toggle','edit-hotzone','deck-viewport'];

  function skipEl(el){
    if (SKIP_TAGS.has(el.tagName)) return true;
    for (const c of SKIP_CLASSES) if (el.classList && el.classList.contains(c)) return true;
    return false;
  }
  function visible(cs){
    if (cs.display === 'none' || cs.visibility === 'hidden') return false;
    if (parseFloat(cs.opacity) === 0) return false;
    return true;
  }
  function isBlockLike(cs){
    const d = cs.display;
    return d === 'block' || d === 'flex' || d === 'grid' ||
           d === 'inline-block' || d === 'list-item' || d === 'flow-root' ||
           d.startsWith('table');
  }
  function hasBgImage(cs){
    return cs.backgroundImage && cs.backgroundImage !== 'none';
  }
  function isGradientText(cs){
    const tf = cs.webkitTextFillColor || cs.textFillColor || '';
    const trans = (v) => v === 'rgba(0, 0, 0, 0)' || v === 'transparent';
    return (trans(tf) || trans(cs.color)) && hasBgImage(cs);
  }
  function borderInfo(cs){
    const w = parseFloat(cs.borderTopWidth) || 0;
    if (w < 0.5) return null;
    if (cs.borderTopStyle === 'none' || cs.borderTopStyle === 'hidden') return null;
    // We only capture uniform borders. If sides differ, first-side approximation.
    return { width: w, color: cs.borderTopColor, style: cs.borderTopStyle };
  }
  function textDescendantBlock(el){
    for (const d of el.querySelectorAll('*')){
      const dcs = getComputedStyle(d);
      if (!visible(dcs)) continue;
      if (!isBlockLike(dcs)) continue;
      let own = '';
      for (const n of d.childNodes){
        if (n.nodeType === 3) own += n.textContent;
      }
      if (own.trim().length > 0) return true;
    }
    return false;
  }
  function commonTextCSS(cs){
    return {
      fontFamily: cs.fontFamily,
      fontSize: parseFloat(cs.fontSize),
      fontWeight: cs.fontWeight,
      fontStyle: cs.fontStyle,
      color: cs.color,
      textAlign: cs.textAlign,
      lineHeight: cs.lineHeight,
      letterSpacing: cs.letterSpacing,
      textTransform: cs.textTransform,
      paddingTop: parseFloat(cs.paddingTop),
      paddingRight: parseFloat(cs.paddingRight),
      paddingBottom: parseFloat(cs.paddingBottom),
      paddingLeft: parseFloat(cs.paddingLeft),
    };
  }
  function extractRuns(el){
    const runs = [];
    function walk(node){
      if (node.nodeType === 3){
        const t = node.textContent;
        if (t.length === 0) return;
        const p = node.parentElement;
        const cs = getComputedStyle(p);
        runs.push({
          text: t,
          fontFamily: cs.fontFamily,
          fontSize: parseFloat(cs.fontSize),
          fontWeight: cs.fontWeight,
          fontStyle: cs.fontStyle,
          color: cs.color,
          textDecoration: cs.textDecorationLine,
          textTransform: cs.textTransform,
          letterSpacing: cs.letterSpacing,
        });
      } else if (node.nodeType === 1){
        if (node.tagName === 'BR'){ runs.push({text:'\n'}); return; }
        for (const c of node.childNodes) walk(c);
      }
    }
    walk(el);
    return runs;
  }

  let capIdx = 0;
  function newCapId(el){
    const id = 's' + slideIdx + '_c' + (capIdx++);
    el.dataset.pptxCap = id;
    return id;
  }

  const slides = Array.from(document.querySelectorAll('.slide'));
  const slide = slides[slideIdx];
  const items = [];
  const bitmapCaps = [];   // {id, hideDescendantText}

  function walk(el){
    if (skipEl(el)) return;
    const cs = getComputedStyle(el);
    if (!visible(cs)) return;
    const r = el.getBoundingClientRect();
    if (r.width < 1 || r.height < 1) return;

    // Gradient-clipped text: capture element itself as one bitmap; do NOT recurse
    if (isGradientText(cs)){
      const id = newCapId(el);
      bitmapCaps.push({id, hideDescendantText: false});
      items.push({type:'bitmap', capId:id, rect:[r.left, r.top, r.width, r.height]});
      return;
    }

    // Raster image
    if (el.tagName === 'IMG'){
      const rad = parseFloat(cs.borderTopLeftRadius) || 0;
      items.push({
        type:'img',
        src: el.currentSrc || el.src,
        rect: [r.left, r.top, r.width, r.height],
        radius: rad,
        objectFit: cs.objectFit,
      });
      return;
    }
    // SVG root (e.g. radar chart) → one bitmap
    if (el.tagName === 'svg' || el.tagName === 'SVG'){
      const id = newCapId(el);
      bitmapCaps.push({id, hideDescendantText: false});
      items.push({type:'bitmap', capId:id, rect:[r.left, r.top, r.width, r.height]});
      return;
    }

    // Background (gradient or color) + border box
    const gradient = hasBgImage(cs);
    const bgColor = cs.backgroundColor;
    const border = borderInfo(cs);
    const radius = parseFloat(cs.borderTopLeftRadius) || 0;
    const boxShadow = cs.boxShadow && cs.boxShadow !== 'none' ? cs.boxShadow : null;

    if (gradient){
      const id = newCapId(el);
      bitmapCaps.push({id, hideDescendantText: false});
      items.push({type:'bitmap', capId:id, rect:[r.left, r.top, r.width, r.height]});
      // Still walk children so text overlays remain editable on top.
    } else {
      const bgRgba = bgColor;
      const bgVisible = bgRgba && bgRgba !== 'rgba(0, 0, 0, 0)' && bgRgba !== 'transparent';
      if (bgVisible || border){
        items.push({
          type:'rect',
          rect: [r.left, r.top, r.width, r.height],
          fill: bgVisible ? bgColor : null,
          radius: radius,
          border: border,
          boxShadow: boxShadow,
        });
      }
    }

    // Text-leaf: block-like element with text, no block-like descendant with own text
    const txt = el.textContent.trim();
    if (txt && isBlockLike(cs) && !textDescendantBlock(el)){
      items.push({
        type:'text',
        rect: [r.left, r.top, r.width, r.height],
        css: commonTextCSS(cs),
        runs: extractRuns(el),
      });
      return;
    }

    for (const c of el.children) walk(c);
  }

  walk(slide);
  return { items, bitmapCaps };
}
"""

ACTIVATE_JS = r"""
(idx) => {
  // Force stage scale = 1
  const stage = document.querySelector('.deck-stage');
  if (stage){ stage.style.transform = 'scale(1)'; }

  const slides = Array.from(document.querySelectorAll('.slide'));
  slides.forEach((s,i)=>{
    s.classList.remove('active','visible');
    if (i === idx){
      s.classList.add('active','visible');
      s.style.visibility='visible';
      s.style.opacity='1';
      s.style.pointerEvents='auto';
      s.style.zIndex='2';
    } else {
      s.style.visibility='hidden';
      s.style.opacity='0';
      s.style.pointerEvents='none';
      s.style.zIndex='';
    }
  });
  // Fast-forward reveal transitions
  document.querySelectorAll('.slide.visible .reveal, .slide.visible .d-stagger > *').forEach(e => {
    e.style.opacity='1';
    e.style.transform='none';
    e.style.transition='none';
    e.style.animation='none';
  });
  return document.querySelectorAll('.slide').length;
}
"""

HIDE_TEXT_JS = r"""
(capId) => {
  const el = document.querySelector('[data-pptx-cap="' + capId + '"]');
  if (!el) return false;
  el.querySelectorAll('*').forEach(c => {
    c.dataset._pOldColor = c.style.color;
    c.dataset._pOldTF = c.style.webkitTextFillColor;
    c.style.color = 'transparent';
    c.style.webkitTextFillColor = 'transparent';
  });
  return true;
}
"""

RESTORE_TEXT_JS = r"""
(capId) => {
  const el = document.querySelector('[data-pptx-cap="' + capId + '"]');
  if (!el) return false;
  el.querySelectorAll('*').forEach(c => {
    c.style.color = c.dataset._pOldColor || '';
    c.style.webkitTextFillColor = c.dataset._pOldTF || '';
    delete c.dataset._pOldColor;
    delete c.dataset._pOldTF;
  });
  return true;
}
"""


# ---------------------------------------------------------------------------
# Extraction driver
# ---------------------------------------------------------------------------
def _log(msg):
    print(msg, flush=True)


def extract_deck(html_path: Path, work_dir: Path):
    work_dir.mkdir(parents=True, exist_ok=True)
    url = "file://" + str(html_path.resolve())
    slides_data: list[dict] = []
    with sync_playwright() as p:
        _log("    · launching chromium")
        browser = p.chromium.launch(headless=True)
        ctx = browser.new_context(
            viewport={"width": SLIDE_W_PX, "height": SLIDE_H_PX},
            device_scale_factor=DPR,
        )
        ctx.set_default_timeout(15000)
        page = ctx.new_page()
        _log("    · loading page")
        page.goto(url, wait_until="load", timeout=30000)
        try:
            page.evaluate("document.fonts && document.fonts.ready")
        except Exception:
            pass
        page.wait_for_timeout(400)

        slide_count = page.evaluate(ACTIVATE_JS, 0)
        _log(f"    → {slide_count} slides detected")
        for i in range(slide_count):
            page.evaluate(ACTIVATE_JS, i)
            page.wait_for_timeout(80)
            data = page.evaluate(EXTRACT_JS, i)
            _log(
                f"    → slide {i+1}/{slide_count}: "
                f"{len(data['items'])} items, {len(data['bitmapCaps'])} bitmaps"
            )
            # Screenshot every bitmap capture region
            for j, cap in enumerate(data["bitmapCaps"]):
                cap_id = cap["id"]
                hide = cap["hideDescendantText"]
                sel = f'[data-pptx-cap="{cap_id}"]'
                if hide:
                    try:
                        page.evaluate(HIDE_TEXT_JS, cap_id)
                    except Exception:
                        pass
                el = page.query_selector(sel)
                if el is None:
                    if hide:
                        try: page.evaluate(RESTORE_TEXT_JS, cap_id)
                        except Exception: pass
                    continue
                png_path = work_dir / f"{cap_id}.png"
                try:
                    el.screenshot(path=str(png_path), timeout=8000)
                except Exception as e:
                    _log(f"        ! bitmap {cap_id} screenshot failed: {e}")
                if hide:
                    try: page.evaluate(RESTORE_TEXT_JS, cap_id)
                    except Exception: pass
            slides_data.append(data)
        browser.close()
    return slides_data


# ---------------------------------------------------------------------------
# PPTX assembly
# ---------------------------------------------------------------------------
def set_shape_fill(shape, rgba):
    if rgba is None:
        shape.fill.background()
        return
    shape.fill.solid()
    shape.fill.fore_color.rgb = to_rgb(rgba)
    if rgba[3] < 1.0:
        # apply alpha via XML
        sp = shape.fill.fore_color._xFill
        # find solidFill/srgbClr and add alpha
        srgb = sp.find(qn("a:srgbClr"))
        if srgb is not None:
            alpha = srgb.find(qn("a:alpha"))
            if alpha is None:
                alpha = etree.SubElement(srgb, qn("a:alpha"))
            alpha.set("val", str(int(rgba[3] * 100000)))


def set_shape_line(shape, border):
    if not border:
        shape.line.fill.background()
        return
    rgba = parse_color(border["color"])
    if rgba is None:
        shape.line.fill.background()
        return
    shape.line.color.rgb = to_rgb(rgba)
    shape.line.width = Emu(px_emu(border["width"]))


def add_rect(slide, item):
    x, y, w, h = item["rect"]
    radius = item.get("radius") or 0
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius > 1 else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(
        shape_type,
        Emu(px_emu(x)),
        Emu(px_emu(y)),
        Emu(px_emu(max(1, w))),
        Emu(px_emu(max(1, h))),
    )
    # Rounded corners: scale adjustment (0..0.5). PPTX adj value = radius / min(w,h) / 2 roughly.
    if radius > 1:
        try:
            adj = min(0.5, radius / max(1, min(w, h)))
            shp.adjustments[0] = adj
        except Exception:
            pass
    set_shape_fill(shp, parse_color(item.get("fill")))
    set_shape_line(shp, item.get("border"))
    shp.shadow.inherit = False
    return shp


def add_bitmap(slide, item, assets_dir: Path, cap_key: str):
    x, y, w, h = item["rect"]
    p = assets_dir / cap_key
    if not p.exists():
        return None
    return slide.shapes.add_picture(
        str(p),
        Emu(px_emu(x)),
        Emu(px_emu(y)),
        Emu(px_emu(max(1, w))),
        Emu(px_emu(max(1, h))),
    )


def add_image(slide, item, base_dir: Path):
    from urllib.parse import unquote, urlparse
    x, y, w, h = item["rect"]
    src = item["src"]
    path = None
    if src.startswith("file://"):
        parsed = urlparse(src)
        path = Path(unquote(parsed.path))
    elif src.startswith("data:"):
        header, _, b64 = src.partition(",")
        import base64 as _b64
        blob = _b64.b64decode(b64)
        import hashlib
        h_name = hashlib.md5(blob).hexdigest()[:10] + ".png"
        tmp = base_dir / f"_inline_{h_name}"
        tmp.write_bytes(blob)
        path = tmp
    else:
        # relative-looking src: resolve against base_dir
        candidate = base_dir / unquote(src)
        if candidate.exists():
            path = candidate
    if not path or not path.exists():
        print(f"    ! missing image on disk: src={src!r} → path={path}", flush=True)
        return None
    # python-pptx only ingests BMP/GIF/JPEG/PNG/TIFF/WMF. Sniff actual bytes
    # (some files use a .png extension but hold WEBP data), and convert if needed.
    try:
        from PIL import Image
        with Image.open(path) as im:
            fmt = (im.format or "").upper()
            if fmt not in {"BMP", "GIF", "JPEG", "PNG", "TIFF", "WMF"}:
                import hashlib
                key = hashlib.md5(str(path).encode()).hexdigest()[:8]
                png_path = base_dir / f"_conv_{path.stem}_{key}.png"
                if not png_path.exists():
                    im.convert("RGBA").save(png_path, "PNG")
                path = png_path
    except Exception as e:
        print(f"    ! image sniff/convert failed for {path.name}: {e}", flush=True)
    return slide.shapes.add_picture(
        str(path),
        Emu(px_emu(x)),
        Emu(px_emu(y)),
        Emu(px_emu(max(1, w))),
        Emu(px_emu(max(1, h))),
    )


def _align_for(css_align: str):
    return {
        "left": PP_ALIGN.LEFT,
        "right": PP_ALIGN.RIGHT,
        "center": PP_ALIGN.CENTER,
        "justify": PP_ALIGN.JUSTIFY,
        "start": PP_ALIGN.LEFT,
        "end": PP_ALIGN.RIGHT,
    }.get(css_align, PP_ALIGN.LEFT)


def add_text(slide, item):
    x, y, w, h = item["rect"]
    css = item["css"]
    tb = slide.shapes.add_textbox(
        Emu(px_emu(x)),
        Emu(px_emu(y)),
        Emu(px_emu(max(1, w))),
        Emu(px_emu(max(1, h))),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(px_emu(css.get("paddingLeft") or 0))
    tf.margin_right = Emu(px_emu(css.get("paddingRight") or 0))
    tf.margin_top = Emu(px_emu(css.get("paddingTop") or 0))
    tf.margin_bottom = Emu(px_emu(css.get("paddingBottom") or 0))
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.auto_size = None

    align = _align_for(css.get("textAlign") or "left")

    # Split runs into paragraphs on "\n"
    runs = item["runs"] or []
    if not runs and item.get("runs") == []:
        runs = [{"text": ""}]

    # Merge into paragraphs
    paragraphs: list[list[dict]] = [[]]
    for r in runs:
        text = r.get("text", "")
        parts = text.split("\n")
        for k, part in enumerate(parts):
            if k > 0:
                paragraphs.append([])
            if part:
                paragraphs[-1].append({**r, "text": part})

    first = True
    for para_runs in paragraphs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        # line spacing from css lineHeight if numeric
        lh = css.get("lineHeight")
        try:
            if lh and lh.endswith("px"):
                lh_px = float(lh[:-2])
                fs = css.get("fontSize") or 16
                if fs > 0:
                    p.line_spacing = round(lh_px / fs, 2)
            elif lh and lh != "normal":
                p.line_spacing = float(lh)
        except Exception:
            pass

        if not para_runs:
            # keep an empty run so paragraph stays
            r = p.add_run()
            r.text = ""
            continue
        for i, run in enumerate(para_runs):
            r = p.add_run()
            text = run.get("text", "")
            # honor text-transform on the run (fallback to css)
            tt = run.get("textTransform") or css.get("textTransform")
            if tt == "uppercase":
                text = text.upper()
            elif tt == "lowercase":
                text = text.lower()
            elif tt == "capitalize":
                text = text.title()
            r.text = text
            font = r.font
            fs = run.get("fontSize") or css.get("fontSize") or 16
            font.size = Pt(round(fs * 0.75, 1))  # px→pt (1pt = 1/72in, 1px = 1/96in → *0.75)
            fam = (run.get("fontFamily") or css.get("fontFamily") or "").split(",")[0].strip().strip('"').strip("'")
            if fam:
                font.name = fam
            weight = run.get("fontWeight") or css.get("fontWeight") or "400"
            try:
                w_num = int(weight)
            except Exception:
                w_num = 700 if weight == "bold" else 400
            font.bold = w_num >= 600
            style = run.get("fontStyle") or css.get("fontStyle") or "normal"
            font.italic = style == "italic"
            color_rgba = parse_color(run.get("color") or css.get("color"))
            if color_rgba:
                font.color.rgb = to_rgb(color_rgba)
            # letter spacing
            ls = run.get("letterSpacing") or css.get("letterSpacing")
            try:
                if ls and ls.endswith("px"):
                    ls_px = float(ls[:-2])
                    # spc unit is 1/100 of a point
                    spc = int(round(ls_px * 0.75 * 100))
                    rPr = r._r.get_or_add_rPr()
                    rPr.set("spc", str(spc))
            except Exception:
                pass
            td = run.get("textDecoration") or ""
            if "underline" in td:
                font.underline = True
    return tb


def build_pptx(slides_data, assets_dir: Path, base_dir: Path, out_path: Path):
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W_PX * EMU_PER_PX)
    prs.slide_height = Emu(SLIDE_H_PX * EMU_PER_PX)
    blank = prs.slide_layouts[6]

    for si, sdata in enumerate(slides_data):
        slide = prs.slides.add_slide(blank)
        # White background — will be covered by rects/bitmaps.
        # Add a full-bleed white rectangle so print-safe background is guaranteed.
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0,
            prs.slide_width, prs.slide_height,
        )
        set_shape_fill(bg, (255, 255, 255, 1.0))
        bg.line.fill.background()

        for item in sdata["items"]:
            t = item["type"]
            try:
                if t == "rect":
                    add_rect(slide, item)
                elif t == "bitmap":
                    cap_key = f"{item['capId']}.png"
                    add_bitmap(slide, item, assets_dir, cap_key)
                elif t == "img":
                    add_image(slide, item, base_dir)
                elif t == "text":
                    add_text(slide, item)
            except Exception as e:
                print(f"  ! slide {si+1} item {t} failed: {e}")
    prs.save(str(out_path))


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------
def main():
    if len(sys.argv) != 3:
        print("usage: html_to_pptx_editable.py <input.html> <output.pptx>")
        sys.exit(2)
    src = Path(sys.argv[1]).resolve()
    out = Path(sys.argv[2]).resolve()
    if not src.exists():
        print(f"input not found: {src}")
        sys.exit(1)
    tmp = Path(tempfile.mkdtemp(prefix="pptx_editable_"))
    print(f"→ extracting DOM to {tmp}")
    slides_data = extract_deck(src, tmp)
    (tmp / "layout.json").write_text(json.dumps(slides_data, indent=1))
    print(f"→ building {out.name}")
    build_pptx(slides_data, tmp, src.parent, out)
    print(f"✓ wrote {out}")


if __name__ == "__main__":
    main()
